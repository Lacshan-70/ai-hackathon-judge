import streamlit as st
from pptx import Presentation
import pdfplumber
import requests
import plotly.graph_objects as go
import numpy as np

# -----------------------------
# PAGE CONFIG
# -----------------------------

st.set_page_config(
    page_title="AI Hackathon Judge",
    page_icon="🤖",
    layout="wide"
)

# -----------------------------
# LOGIN SYSTEM
# -----------------------------

if "logged_in" not in st.session_state:
    st.session_state.logged_in = False

def login():

    st.title("🤖 AI Hackathon Judge Login")

    username = st.text_input("Username")
    password = st.text_input("Password", type="password")

    if st.button("Login"):

        if username == "admin" and password == "admin123":
            st.session_state.logged_in = True
            st.rerun()

        else:
            st.error("Invalid credentials")

if not st.session_state.logged_in:
    login()
    st.stop()

# -----------------------------
# HEADER
# -----------------------------

st.title("🤖 AI Hackathon PPT Judge")
st.write("Upload your presentation and let AI evaluate it like a real jury.")

# -----------------------------
# FILE UPLOADER
# -----------------------------

uploaded_file = st.file_uploader(
    "Upload PPT or PDF",
    type=["pptx", "pdf"]
)

# -----------------------------
# PPT TEXT EXTRACTION
# -----------------------------

def extract_ppt_text(file):

    prs = Presentation(file)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text

# -----------------------------
# PDF TEXT EXTRACTION
# -----------------------------

def extract_pdf_text(file):

    text = ""

    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:

            page_text = page.extract_text()

            if page_text:
                text += page_text

    return text

# -----------------------------
# AI EVALUATION (GROQ)
# -----------------------------

def evaluate_presentation(text):

    url = "https://api.groq.com/openai/v1/chat/completions"

    headers = {
        "Authorization": f"Bearer {st.secrets['GROQ_API_KEY']}",
        "Content-Type": "application/json"
    }

    prompt = f"""
You are an expert hackathon judge.

Evaluate this presentation.

Return scores out of 10 in this exact format:

Problem: X
Innovation: X
Feasibility: X
Impact: X
Implementation: X
Presentation: X

Then give bullet point feedback.

Presentation:
{text[:4000]}
"""

    payload = {
        "model": "llama3-70b-8192",
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.3
    }

    response = requests.post(url, headers=headers, json=payload)

    if response.status_code != 200:
        return None, f"API Error: {response.text}"

    result = response.json()

    if "choices" in result:
        return result["choices"][0]["message"]["content"], None

    return None, str(result)

# -----------------------------
# SCORE EXTRACTION
# -----------------------------

def extract_scores(text):

    scores = {
        "Problem":5,
        "Innovation":5,
        "Feasibility":5,
        "Impact":5,
        "Implementation":5,
        "Presentation":5
    }

    for key in scores:
        if key in text:
            try:
                scores[key] = int(text.split(key+":")[1].split()[0])
            except:
                pass

    return scores

# -----------------------------
# PROCESS FILE
# -----------------------------

if uploaded_file:

    if uploaded_file.name.endswith(".pptx"):
        full_text = extract_ppt_text(uploaded_file)

    elif uploaded_file.name.endswith(".pdf"):
        full_text = extract_pdf_text(uploaded_file)

    st.subheader("Extracted Content")
    st.write(full_text[:1000])

    st.write("---")

    with st.spinner("AI is evaluating your presentation..."):

        feedback, error = evaluate_presentation(full_text)

    if error:
        st.error(error)
        st.stop()

    scores = extract_scores(feedback)

    problem = scores["Problem"]
    innovation = scores["Innovation"]
    feasibility = scores["Feasibility"]
    impact = scores["Impact"]
    implementation = scores["Implementation"]
    presentation = scores["Presentation"]

# -----------------------------
# AI FEEDBACK
# -----------------------------

    st.subheader("AI Jury Feedback")
    st.write(feedback)

# -----------------------------
# SCORE CALCULATION
# -----------------------------

    total = (
        problem +
        innovation +
        feasibility +
        impact +
        implementation +
        presentation
    )

    probability = int((total / 60) * 100)

# -----------------------------
# SCORE CARDS
# -----------------------------

    st.write("---")
    st.subheader("Score Summary")

    col1, col2, col3 = st.columns(3)

    with col1:
        st.metric("Problem", f"{problem}/10")
        st.metric("Impact", f"{impact}/10")

    with col2:
        st.metric("Innovation", f"{innovation}/10")
        st.metric("Implementation", f"{implementation}/10")

    with col3:
        st.metric("Feasibility", f"{feasibility}/10")
        st.metric("Presentation", f"{presentation}/10")

# -----------------------------
# SELECTION PROBABILITY
# -----------------------------

    st.write("---")

    st.subheader("Selection Probability")

    st.progress(probability / 100)

    st.write(f"### {probability}% Chance of Selection")

    if probability > 75:
        st.success("🔥 High Chance of Selection")

    elif probability > 50:
        st.warning("⚠ Moderate Chance")

    else:
        st.error("❌ Low Chance")

# -----------------------------
# RADAR CHART
# -----------------------------

    st.write("---")
    st.subheader("Evaluation Radar Chart")

    labels = [
        "Problem",
        "Innovation",
        "Feasibility",
        "Impact",
        "Implementation",
        "Presentation"
    ]

    values = [
        problem,
        innovation,
        feasibility,
        impact,
        implementation,
        presentation
    ]

    values += values[:1]

    angles = np.linspace(0, 2*np.pi, len(labels), endpoint=False).tolist()
    angles += angles[:1]

    fig = go.Figure()

    fig.add_trace(go.Scatterpolar(
        r=values,
        theta=labels + [labels[0]],
        fill="toself"
    ))

    fig.update_layout(
        polar=dict(radialaxis=dict(range=[0,10])),
        showlegend=False
    )

    st.plotly_chart(fig, use_container_width=True)